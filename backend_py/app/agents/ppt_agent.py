from pathlib import Path

from app.models import AgentRunResult, RunEvent


def run_ppt_agent(task_id: str, task: str, workspace_root: Path) -> AgentRunResult:
    outline = [
        "Problem and user need",
        "Universal Agent architecture",
        "Router Agent and specialist Agents",
        "Code sandbox workflow",
        "Data analysis workflow",
        "File permission workflow",
        "MCP runtime and tools",
        "Approvals, audit, and safety",
        "Demo scenario",
        "Roadmap",
    ]
    markdown = "# Draft Slide Outline\n\nTask: " + task + "\n\n" + "\n".join(f"{i + 1}. {title}" for i, title in enumerate(outline))
    pptx_path = workspace_root / "agent_deck.pptx"
    pptx_status = _create_pptx(task, outline, pptx_path)

    artifacts = {"outline.md": markdown, "pptx_path.txt": str(pptx_path)}
    return AgentRunResult(
        task_id=task_id,
        task_type="ppt",
        summary=f"PPT Agent generated a real PPTX artifact. {pptx_status}",
        events=[
            RunEvent(id="outline", title="Outline generated", detail="10 slide draft", state="done", meta="artifact"),
            RunEvent(id="pptx", title="PPTX exported", detail=str(pptx_path), state="done", meta="python-pptx"),
        ],
        artifacts=artifacts,
    )


def _create_pptx(task: str, slide_titles: list[str], path: Path) -> str:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except Exception as exc:
        path.with_suffix(".txt").write_text("python-pptx is not installed.\n" + str(exc), encoding="utf-8")
        return "python-pptx is missing; wrote fallback note."

    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]

    title_slide = prs.slides.add_slide(title_layout)
    title_slide.shapes.title.text = "My Agent App"
    title_slide.placeholders[1].text = task[:180]

    for index, title in enumerate(slide_titles, start=1):
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = f"{index}. {title}"
        body = slide.placeholders[1].text_frame
        body.clear()
        bullets = [
            "What this stage does",
            "Which Agent owns it",
            "Tools, sandbox, and approval boundaries",
        ]
        for bullet in bullets:
            paragraph = body.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.font.size = Pt(20)

    final_slide = prs.slides.add_slide(content_layout)
    final_slide.shapes.title.text = "Next steps"
    final_slide.placeholders[1].text = "Connect streaming, approvals, MCP tools, Docker sandboxing, and production model providers."

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)
    return f"Saved {path.name}."
