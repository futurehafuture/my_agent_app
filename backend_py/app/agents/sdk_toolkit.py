import json
import os
import shutil
import subprocess
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from app.runtime.permissions import ensure_within_root, is_risky_command
from app.tools.browser_tools import web_search
from app.tools.diff_tools import directory_diff, git_diff
from app.tools.shell_tools import suggest_validation_commands

IGNORE_DIRS = {".git", "node_modules", ".next", "dist", "build", "out", ".venv", "__pycache__"}


@dataclass
class AgentWorkspace:
    root: Path
    source_project: Path | None = None
    repo: Path | None = None
    data_source: Path | None = None
    allowed_folder: Path | None = None
    artifacts: dict[str, str] = field(default_factory=dict)
    tool_logs: list[str] = field(default_factory=list)

    def log(self, message: str) -> None:
        self.tool_logs.append(message[:4000])


def prepare_code_workspace(workspace: AgentWorkspace, project_path: str | None) -> None:
    if not project_path:
        return
    source = Path(project_path).expanduser().resolve()
    if not source.exists() or not source.is_dir():
        raise FileNotFoundError(f"Project folder does not exist: {source}")
    repo = workspace.root / "repo"
    if repo.exists():
        shutil.rmtree(repo)

    def ignore(_: str, names: list[str]) -> set[str]:
        return {name for name in names if name in IGNORE_DIRS or name.endswith(".log")}

    shutil.copytree(source, repo, ignore=ignore)
    workspace.source_project = source
    workspace.repo = repo


def prepare_data_workspace(workspace: AgentWorkspace, data_path: str | None) -> None:
    if not data_path:
        return
    source = Path(data_path).expanduser().resolve()
    if not source.exists():
        raise FileNotFoundError(f"Data path does not exist: {source}")
    data_dir = workspace.root / "data"
    data_dir.mkdir(parents=True, exist_ok=True)
    if source.is_file():
        shutil.copy2(source, data_dir / source.name)
    elif source.is_dir():
        for child in source.iterdir():
            if child.is_file() and child.suffix.lower() in {".csv", ".xlsx", ".xls"}:
                shutil.copy2(child, data_dir / child.name)
    workspace.data_source = source


def prepare_file_workspace(workspace: AgentWorkspace, allowed_folder: str | None) -> None:
    if allowed_folder:
        folder = Path(allowed_folder).expanduser().resolve()
        if not folder.exists() or not folder.is_dir():
            raise FileNotFoundError(f"Allowed folder does not exist: {folder}")
        workspace.allowed_folder = folder


def create_tools(workspace: AgentWorkspace) -> list[Any]:
    from agents import function_tool

    @function_tool
    def list_files(path: str = ".") -> str:
        """List files under the sandbox repo or workspace. Path must be relative."""
        base = workspace.repo or workspace.root
        target = ensure_within_root(base, base / path)
        if not target.exists():
            return "Path not found."
        result = sorted(item.name + ("/" if item.is_dir() else "") for item in target.iterdir())[:300]
        workspace.log(f"list_files({path}) -> {len(result)} items")
        return "\n".join(result)

    @function_tool
    def read_file(path: str) -> str:
        """Read a UTF-8 text file from the sandbox repo. Path must be relative."""
        base = workspace.repo or workspace.root
        target = ensure_within_root(base, base / path)
        text = target.read_text(encoding="utf-8", errors="replace")[:50000]
        workspace.log(f"read_file({path}) -> {len(text)} chars")
        return text

    @function_tool
    def write_file(path: str, content: str) -> str:
        """Write a UTF-8 text file inside the sandbox repo. Path must be relative."""
        base = workspace.repo or workspace.root
        target = ensure_within_root(base, base / path)
        target.parent.mkdir(parents=True, exist_ok=True)
        target.write_text(content, encoding="utf-8")
        workspace.log(f"write_file({path}) -> {len(content)} chars")
        return f"Wrote {path}"

    @function_tool
    def run_command(command: list[str], timeout_seconds: int = 90) -> str:
        """Run a safe command in the sandbox repo. Risky commands are blocked."""
        base = workspace.repo or workspace.root
        joined = " ".join(command)
        if is_risky_command(joined):
            workspace.log(f"run_command blocked: {joined}")
            return f"BLOCKED: command requires approval: {joined}"
        completed = subprocess.run(
            command,
            cwd=base,
            text=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            timeout=min(timeout_seconds, 180),
            check=False,
        )
        output = f"exit={completed.returncode}\nSTDOUT:\n{completed.stdout[-20000:]}\nSTDERR:\n{completed.stderr[-20000:]}"
        workspace.log(f"run_command({joined}) -> exit {completed.returncode}")
        return output

    @function_tool
    def get_diff() -> str:
        """Return the diff between the original project and sandbox repo."""
        if workspace.repo and workspace.source_project:
            diff = git_diff(workspace.repo) or directory_diff(workspace.source_project, workspace.repo)
        else:
            diff = "No code workspace available."
        workspace.artifacts["diff.patch"] = diff
        workspace.log("get_diff()")
        return diff[:50000]

    @function_tool
    def suggest_validation() -> str:
        """Suggest likely validation commands for the sandbox repo."""
        if not workspace.repo:
            return "No repo workspace available."
        commands = suggest_validation_commands(workspace.repo)
        result = "\n".join(commands) or "No obvious validation command detected."
        workspace.artifacts["validation_commands.txt"] = result
        workspace.log("suggest_validation()")
        return result

    @function_tool
    def summarize_data() -> str:
        """Summarize CSV/XLSX files copied into the data workspace."""
        data_dir = workspace.root / "data"
        try:
            import pandas as pd
        except Exception as exc:
            return f"pandas is not installed: {exc}"
        sections = ["# Data Summary"]
        files = [p for p in data_dir.glob("*") if p.suffix.lower() in {".csv", ".xlsx", ".xls"}]
        for file in files:
            try:
                df = pd.read_csv(file) if file.suffix.lower() == ".csv" else pd.read_excel(file)
                sections.append(f"\n## {file.name}\nRows: {len(df)}\nColumns: {len(df.columns)}\nColumn names: {', '.join(map(str, df.columns[:30]))}")
                numeric = df.select_dtypes(include="number")
                if not numeric.empty:
                    sections.append("\nNumeric summary:\n" + numeric.describe().to_markdown())
            except Exception as exc:
                sections.append(f"\n## {file.name}\nFailed: {exc}")
        report = "\n".join(sections) if files else "No CSV/XLSX files found."
        workspace.artifacts["data_report.md"] = report
        workspace.log(f"summarize_data() -> {len(files)} files")
        return report[:50000]

    @function_tool
    def create_pptx(title: str, slide_titles: list[str]) -> str:
        """Create a PPTX deck artifact from slide titles."""
        path = workspace.root / "agent_deck.pptx"
        try:
            from pptx import Presentation
            from pptx.util import Pt
        except Exception as exc:
            return f"python-pptx is not installed: {exc}"
        prs = Presentation()
        cover = prs.slides.add_slide(prs.slide_layouts[0])
        cover.shapes.title.text = title
        cover.placeholders[1].text = "Generated by My Agent App"
        for index, slide_title in enumerate(slide_titles[:20], start=1):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"{index}. {slide_title}"
            body = slide.placeholders[1].text_frame
            body.clear()
            for bullet in ["Key point", "Evidence or detail", "Next action"]:
                p = body.add_paragraph()
                p.text = bullet
                p.font.size = Pt(20)
        prs.save(path)
        workspace.artifacts["pptx_path.txt"] = str(path)
        workspace.log(f"create_pptx({title}) -> {path}")
        return f"Saved PPTX to {path}"

    @function_tool
    def browser_search(query: str) -> str:
        """Search the web and return result cards with titles, URLs, and snippets."""
        results = web_search(query, max_results=6)
        text = json.dumps(results, ensure_ascii=False, indent=2)
        workspace.artifacts["research_results.json"] = text
        workspace.log(f"browser_search({query}) -> {len(results)} results")
        return text

    @function_tool
    def scan_allowed_folder() -> str:
        """Scan the user-approved folder and propose a file organization plan. Does not move files."""
        if not workspace.allowed_folder:
            return "No approved folder configured."
        groups: dict[str, list[str]] = {}
        for child in workspace.allowed_folder.iterdir():
            if child.is_file():
                groups.setdefault(child.suffix.lower() or "no-extension", []).append(child.name)
        lines = [f"# File plan for {workspace.allowed_folder}"]
        for suffix, names in sorted(groups.items()):
            lines.append(f"- {len(names)} {suffix} file(s): move to {suffix.replace('.', '').upper() or 'MISC'}/ after approval")
        plan = "\n".join(lines) if len(lines) > 1 else "No files found."
        workspace.artifacts["file_plan.md"] = plan
        workspace.log("scan_allowed_folder()")
        return plan

    return [
        list_files,
        read_file,
        write_file,
        run_command,
        get_diff,
        suggest_validation,
        summarize_data,
        create_pptx,
        browser_search,
        scan_allowed_folder,
    ]
